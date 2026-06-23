from __future__ import annotations

from io import BytesIO
import re

from app.services.limits import LimitViolation


IMAGE_SUFFIXES = {"png", "jpg", "jpeg", "webp", "tif", "tiff"}
TEXT_SUFFIXES = {"txt", "md", "markdown", "csv"}
SUPPORTED_SUFFIXES = TEXT_SUFFIXES | IMAGE_SUFFIXES | {"pdf", "docx", "pptx"}
UNSUPPORTED_FILE_TYPE_MESSAGE = (
    "Unsupported file type. Upload PDF, DOCX, PPTX, TXT, Markdown, CSV, PNG, JPG, WEBP, or TIFF."
)

STOPWORDS = {
    "about",
    "after",
    "again",
    "also",
    "and",
    "are",
    "because",
    "been",
    "before",
    "being",
    "between",
    "chunk",
    "document",
    "from",
    "have",
    "into",
    "image",
    "notes",
    "ocr",
    "page",
    "slide",
    "that",
    "the",
    "their",
    "then",
    "there",
    "these",
    "this",
    "through",
    "with",
    "would",
}


def is_supported_filename(filename: str) -> bool:
    return _file_suffix(filename) in SUPPORTED_SUFFIXES


def extract_text(
    filename: str,
    content: bytes,
    ocr_client=None,
    max_pdf_pages: int | None = None,
    max_ocr_pages: int | None = None,
) -> str:
    suffix = _file_suffix(filename)
    if suffix not in SUPPORTED_SUFFIXES:
        raise RuntimeError(UNSUPPORTED_FILE_TYPE_MESSAGE)

    if suffix == "pdf":
        return _extract_pdf_text(filename, content, ocr_client, max_pdf_pages, max_ocr_pages)

    if suffix == "docx":
        return _extract_docx_text(content)

    if suffix == "pptx":
        return _extract_pptx_text(content)

    if suffix in IMAGE_SUFFIXES:
        return _extract_image_text(filename, content, ocr_client)

    return _tag_section("document", content.decode("utf-8", errors="ignore").strip())


def chunk_text(text: str, chunk_size: int = 900, overlap: int = 160) -> list[dict]:
    cleaned = re.sub(r"\n{3,}", "\n\n", text).strip()
    if not cleaned:
        return []

    sections = _sections_from_text(cleaned)
    chunked: list[dict] = []
    chunk_index = 1
    for label, body in sections:
        for chunk in _split_section_body(body, chunk_size, overlap):
            text_value = f"[{label}]\n{chunk.strip()}" if label != "document" else chunk.strip()
            chunked.append(
                {
                    "text": text_value,
                    "section": f"{label}, chunk {chunk_index}" if label != "document" else f"chunk {chunk_index}",
                }
            )
            chunk_index += 1
    return chunked


def _split_section_body(text: str, chunk_size: int, overlap: int) -> list[str]:
    try:
        from langchain_text_splitters import RecursiveCharacterTextSplitter

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=overlap,
            separators=["\n\n", "\n", ". ", " ", ""],
        )
        return splitter.split_text(text.strip())
    except ModuleNotFoundError:
        return _fallback_split(text.strip(), chunk_size, overlap)


def _extract_pdf_text(
    filename: str,
    content: bytes,
    ocr_client=None,
    max_pdf_pages: int | None = None,
    max_ocr_pages: int | None = None,
) -> str:
    try:
        from pypdf import PdfReader
    except ModuleNotFoundError as exc:
        raise RuntimeError("PDF extraction requires pypdf. Install backend requirements.") from exc

    reader = PdfReader(BytesIO(content))
    if max_pdf_pages and len(reader.pages) > max_pdf_pages:
        raise LimitViolation(413, f"PDF has {len(reader.pages)} pages, above the {max_pdf_pages} page beta limit.")

    pages: list[str] = []
    blank_pages: list[int] = []
    ocr_pages = 0
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(_tag_section(f"page {index}", text.strip()))
            continue

        blank_pages.append(index)
        if ocr_client:
            if max_ocr_pages and ocr_pages >= max_ocr_pages:
                raise LimitViolation(413, f"PDF needs OCR for more than {max_ocr_pages} pages.")
            ocr_pages += 1
            image = _render_pdf_page(content, index)
            ocr_text, _ = ocr_client.ocr_image(f"{filename}-page-{index}.png", image)
            if ocr_text.strip():
                pages.append(_tag_section(f"page {index} OCR", ocr_text.strip()))

    if not pages and blank_pages:
        raise RuntimeError("This PDF has no selectable text. OCR needs OPENAI_API_KEY and a vision-capable OPENAI_MODEL.")

    return "\n\n".join(pages).strip()


def _file_suffix(filename: str) -> str:
    return filename.lower().rsplit(".", 1)[-1] if "." in filename else "txt"


def _extract_docx_text(content: bytes) -> str:
    try:
        from docx import Document
    except ModuleNotFoundError as exc:
        raise RuntimeError("DOCX extraction requires python-docx. Install backend requirements.") from exc

    document = Document(BytesIO(content))
    parts = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return _tag_section("document", "\n".join(parts).strip())


def _extract_pptx_text(content: bytes) -> str:
    try:
        from pptx import Presentation
    except ModuleNotFoundError as exc:
        raise RuntimeError("Slide extraction requires python-pptx. Install backend requirements.") from exc

    presentation = Presentation(BytesIO(content))
    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            slides.append(_tag_section(f"slide {index}", "\n".join(parts)))
    return "\n\n".join(slides).strip()


def _extract_image_text(filename: str, content: bytes, ocr_client=None) -> str:
    if not ocr_client:
        raise RuntimeError("Image OCR needs OPENAI_API_KEY and a vision-capable OPENAI_MODEL.")
    text, _ = ocr_client.ocr_image(filename, content)
    if not text.strip():
        raise RuntimeError("No readable text found in this image.")
    return _tag_section("image OCR", text.strip())


def _render_pdf_page(content: bytes, page_number: int) -> bytes:
    try:
        import fitz
    except ModuleNotFoundError as exc:
        raise RuntimeError("Scanned PDF OCR requires PyMuPDF. Install backend requirements.") from exc

    with fitz.open(stream=content, filetype="pdf") as document:
        page = document.load_page(page_number - 1)
        pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
        return pixmap.tobytes("png")


def _fallback_split(text: str, chunk_size: int, overlap: int) -> list[str]:
    paragraphs = [part.strip() for part in text.split("\n\n") if part.strip()]
    chunks: list[str] = []
    current = ""

    for paragraph in paragraphs:
        next_value = f"{current}\n\n{paragraph}".strip() if current else paragraph
        if len(next_value) <= chunk_size:
            current = next_value
            continue

        if current:
            chunks.append(current)
        current = paragraph

        while len(current) > chunk_size:
            chunks.append(current[:chunk_size])
            current = current[max(0, chunk_size - overlap) :]

    if current:
        chunks.append(current)

    return chunks


def _sections_from_text(text: str) -> list[tuple[str, str]]:
    sections: list[tuple[str, str]] = []
    current_label = "document"
    current_lines: list[str] = []
    marker_pattern = re.compile(r"^\[(page \d+(?: OCR)?|slide \d+|image OCR|document)\]\s*$", re.IGNORECASE)

    for line in text.splitlines():
        marker = marker_pattern.match(line.strip())
        if marker:
            if "\n".join(current_lines).strip():
                sections.append((current_label, "\n".join(current_lines).strip()))
            current_label = marker.group(1)
            current_lines = []
            continue
        current_lines.append(line)

    if "\n".join(current_lines).strip():
        sections.append((current_label, "\n".join(current_lines).strip()))

    return sections or [("document", text)]


def _tag_section(label: str, text: str) -> str:
    return f"[{label}]\n{text.strip()}" if text.strip() else ""


def extract_keywords(text: str, limit: int = 8) -> list[str]:
    words = re.findall(r"[A-Za-z][A-Za-z\-]{3,}", text.lower())
    counts: dict[str, int] = {}
    for word in words:
        if word in STOPWORDS:
            continue
        counts[word] = counts.get(word, 0) + 1
    ranked = sorted(counts.items(), key=lambda item: (-item[1], item[0]))
    return [word for word, _ in ranked[:limit]]
